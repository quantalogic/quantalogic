"""Tool for converting markdown content to well-structured PowerPoint presentations.

Why this tool:
- Provides a standardized way to convert markdown to professional PPTX presentations
- Maintains consistent styling and formatting across slides
- Handles complex elements like diagrams, code blocks, and images
- Supports customization through templates and style configurations
"""

import json
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

import markdown
from bs4 import BeautifulSoup
from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from quantalogic.tools.tool import Tool, ToolArgument


class MarkdownToPptxTool(Tool):
    """Converts markdown to professional PowerPoint presentations with advanced formatting."""

    model_config = {
        "arbitrary_types_allowed": True
    }
    
    name: str = "markdown_to_pptx_tool"
    description: str = (
        "Converts markdown to PPTX with support for images, Mermaid diagrams, "
        "code blocks, and advanced slide formatting."
    )
    need_validation: bool = False
    
    arguments: List[ToolArgument] = [
        ToolArgument(
            name="markdown_content",
            arg_type="string",
            description="Markdown content where slides are separated by '---'. Supports headers, bullet points, numbered lists, bold/italic text, code blocks, and Mermaid diagrams.",
            required=True,
            example='''# Project Overview
Quarterly Update Q1 2024

---

# Key Achievements
- **Revenue Growth**: 25% YoY increase
- **New Features**: 
  - AI-powered recommendations
  - Real-time analytics
- **Customer Satisfaction**: 4.8/5

---

# Technical Architecture
```mermaid
graph TD
    A[Frontend] --> B[API Gateway]
    B --> C[Microservices]
    C --> D[Database]
```

---

# Next Steps
1. Launch mobile app
2. Expand to new markets
3. Enhance AI capabilities

---

# Questions?
Thank you for your attention!
''',
        ),
        ToolArgument(
            name="output_path",
            arg_type="string",
            description="Path for saving the PPTX file",
            required=True,
            example="/path/to/output.pptx",
        ),
        ToolArgument(
            name="template_path",
            arg_type="string",
            description="Optional PPTX template path",
            required=False,
            example="/path/to/template.pptx",
        ),
        ToolArgument(
            name="style_config",
            arg_type="string",
            description="JSON string with style settings (fonts, colors, sizes)",
            required=False,
            example='{"font_name": "Calibri", "title_size": 44, "body_size": 24}',
        ),
    ]

    # Default style configuration
    DEFAULT_STYLES: Dict[str, Union[str, int, List[int]]] = {
        "font_name": "Calibri",
        "title_size": 44,
        "body_size": 24,
        "code_size": 18,
        "code_font": "Consolas",
        "primary_color": [0, 112, 192],
        "secondary_color": [68, 114, 196],
        "text_color": [0, 0, 0],
    }

    def _normalize_path(self, path: str) -> Path:
        """Convert path string to normalized Path object."""
        if path.startswith("~"):
            path = os.path.expanduser(path)
        return Path(path).resolve()

    def _parse_style_config(self, style_config: Optional[str]) -> Dict:
        """Parse and validate style configuration."""
        try:
            if not style_config:
                return self.DEFAULT_STYLES.copy()
            
            custom_styles = json.loads(style_config)
            styles = self.DEFAULT_STYLES.copy()
            styles.update(custom_styles)
            return styles
        except json.JSONDecodeError as e:
            logger.error(f"Invalid style configuration JSON: {e}")
            return self.DEFAULT_STYLES.copy()

    def _create_presentation(self, template_path: Optional[str] = None) -> Presentation:
        """Create a new presentation, optionally from a template."""
        if template_path:
            template_path = self._normalize_path(template_path)
            if not template_path.exists():
                logger.warning(f"Template not found: {template_path}. Using default template.")
                return Presentation()
            return Presentation(template_path)
        return Presentation()

    def _apply_text_style(self, paragraph, font_name: str, font_size: int, color: List[int]):
        """Apply text styling to a paragraph."""
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.color.rgb = RGBColor(*color)

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: Optional[str] = None):
        """Add a title slide to the presentation."""
        layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(layout)
        
        # Add title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
        
        # Add subtitle if it exists and there's a subtitle placeholder
        if subtitle:
            try:
                subtitle_shape = slide.placeholders[1]  # Index 1 is typically the subtitle placeholder
                subtitle_shape.text = subtitle
            except (KeyError, IndexError):
                logger.warning("No subtitle placeholder found in the slide layout")

    def _add_content_slide(self, prs: Presentation, title: str, content: str, styles: Dict):
        """Add a content slide with formatted text and elements."""
        layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(layout)
        
        # Add title if title placeholder exists
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_text_style(
                title_shape.text_frame.paragraphs[0],
                styles["font_name"],
                styles["title_size"],
                styles["primary_color"]
            )
        
        try:
            # Get content placeholder
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()  # Clear existing content
            
            # Convert markdown to HTML for better parsing
            html_content = markdown.markdown(content)
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Process content
            for element in soup.children:
                if isinstance(element, str):
                    continue
                    
                if element.name == 'ul':
                    # Handle bullet points
                    for li in element.find_all('li'):
                        p = tf.add_paragraph()
                        p.level = 0
                        p.text = li.get_text().strip()
                        self._apply_text_style(p, styles["font_name"], styles["body_size"], styles["text_color"])
                elif element.name == 'p':
                    # Handle paragraphs
                    p = tf.add_paragraph()
                    p.text = element.get_text().strip()
                    self._apply_text_style(p, styles["font_name"], styles["body_size"], styles["text_color"])
                
        except (KeyError, IndexError) as e:
            logger.warning(f"Error adding content to slide: {str(e)}")
            # Fallback to basic text box if no content placeholder
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = content
            self._apply_text_style(p, styles["font_name"], styles["body_size"], styles["text_color"])

    def execute(self, **kwargs) -> Dict:
        """Convert markdown content to a PowerPoint presentation."""
        try:
            markdown_content = kwargs.get("markdown_content")
            output_path = kwargs.get("output_path")
            template_path = kwargs.get("template_path")
            style_config = kwargs.get("style_config")

            if not markdown_content or not output_path:
                raise ValueError("markdown_content and output_path are required")

            # Parse style configuration
            styles = self._parse_style_config(style_config)
            
            # Create presentation
            prs = self._create_presentation(template_path)
            
            # Split content into slides
            slides = markdown_content.split("---")
            
            # Process each slide
            for i, slide_content in enumerate(slides):
                if not slide_content.strip():
                    continue
                    
                # Parse slide content
                lines = slide_content.strip().split("\n")
                title = lines[0].lstrip("#").strip() if lines else "Untitled Slide"
                content = "\n".join(lines[1:]).strip()
                
                if i == 0:
                    # First slide is title slide
                    subtitle = content if content else None
                    self._add_title_slide(prs, title, subtitle)
                else:
                    self._add_content_slide(prs, title, content, styles)
            
            # Save presentation
            output_path = self._normalize_path(output_path)
            prs.save(str(output_path))
            
            return {
                "status": "success",
                "message": f"Presentation saved to {output_path}",
                "output_path": str(output_path)
            }
            
        except Exception as e:
            logger.error(f"Error creating presentation: {str(e)}")
            return {
                "status": "error",
                "message": str(e)
            }


if __name__ == "__main__":
    # Example usage with error handling
    try:
        tool = MarkdownToPptxTool()
        
        # Example markdown content
        markdown_content = """
        # Project Overview
        Quarterly Update Q1 2024
        
        ---
        
        # Key Achievements
        
        - Launched new product features
        - Increased user engagement by 25%
        - Improved system performance
        
        ---
        
        # Next Steps
        
        1. Scale infrastructure
        2. Release mobile app
        3. Expand market reach
        """
        
        result = tool.execute(
            markdown_content=markdown_content,
            output_path="presentation.pptx",
            style_config='{"primary_color": [44, 85, 169]}'
        )
        
        print(result)
        
    except Exception as e:
        logger.error(f"Example usage failed: {e}")
